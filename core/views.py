from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status, permissions
from rest_framework_simplejwt.tokens import RefreshToken
from .serializers import RegisterSerializer, LoginSerializer, LessonSerializer
from .models import Lesson
from rest_framework.parsers import MultiPartParser
import os
from django.core.files.storage import default_storage
from django.conf import settings
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from core.models import User, Lesson




class RegisterView(APIView):
    def post(self, request):
        email = request.data.get("email")
        password = request.data.get("password")
        role = request.data.get("role", "student")

        if not email or not password:
            return Response({"error": "Email and password required"}, status=status.HTTP_400_BAD_REQUEST)

        # Try creating the user
        try:
            user = User.objects.create_user(email=email, password=password, role=role)
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_400_BAD_REQUEST)

        return Response({"email": user.email, "role": user.role}, status=status.HTTP_201_CREATED)





class LoginView(APIView):
    def post(self, request):
        serializer = LoginSerializer(data=request.data)
        if serializer.is_valid():
            user = serializer.validated_data["user"]
            refresh = RefreshToken.for_user(user)
            return Response({
                "refresh": str(refresh),
                "access": str(refresh.access_token)
            }, status=status.HTTP_200_OK)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)



class LessonView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def get(self, request):
        lessons = Lesson.objects.all()
        serializer = LessonSerializer(lessons, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    def post(self, request):
        serializer = LessonSerializer(data=request.data)
        if serializer.is_valid():
            serializer.save()
            return Response(serializer.data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)



#UPLOAD_DIR = "uploads"
#os.makedirs(UPLOAD_DIR, exist_ok=True)

#class FileUploadView(APIView):
#    parser_classes = [MultiPartParser]

#    def post(self, request):
 #       file_obj = request.FILES.get('file')
  #      if not file_obj:
   #         return Response({"error": "No file provided"}, status=status.HTTP_400_BAD_REQUEST)
#
 #       raw_text = file_obj.read().decode('utf-8')
#
 #       file_path = os.path.join(UPLOAD_DIR, file_obj.name)
  #      with open(file_path, 'wb') as f:
   #         f.write(file_obj.read())
#
 #       return Response({
  #          "message": "File uploaded successfully",
   #         "filename": file_obj.name,
    #        "content": raw_text
     #   }, status=status.HTTP_201_CREATED)



UPLOAD_DIR = os.path.join(settings.BASE_DIR, "uploads")
EXPORT_DIR = os.path.join(settings.BASE_DIR, "exports")
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(EXPORT_DIR, exist_ok=True)


cv_results = {}
uploaded_texts = {}

import os
from django.utils.crypto import get_random_string
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status, permissions

class UploadFileView(APIView):
    permission_classes = [permissions.IsAuthenticated]  # JWT required

    def post(self, request):
        file = request.FILES.get('file')
        if not file:
            return Response({"error": "No file provided"}, status=status.HTTP_400_BAD_REQUEST)

        upload_dir = "uploaded_files"
        os.makedirs(upload_dir, exist_ok=True)

    
        filename, ext = os.path.splitext(file.name)
        new_filename = f"{filename}_{get_random_string(6)}{ext}"

        file_path = os.path.join(upload_dir, new_filename)
        with open(file_path, 'wb+') as destination:
            for chunk in file.chunks():
                destination.write(chunk)

        return Response({"message": f"File '{new_filename}' uploaded successfully"}, status=201)



class AnalyzeCVView(APIView):
    def post(self, request):
        file = request.FILES.get('file')
        if not file:
            return Response({"detail": "No file uploaded"}, status=400)
        
        content = file.read().decode("utf-8")
        lines = content.split("\n")
        skills = [line for line in lines if any(skill in line for skill in ["Python","Java","SQL"])]
        experience_years = sum([int(s) for s in lines if s.isdigit()])

        cv_result = {"filename": file.name, "skills_found": skills, "experience_years": experience_years}
        cv_results[file.name] = cv_result

        return Response(cv_result)


class RoadmapView(APIView):
    def get(self, request, filename):
        if filename not in cv_results:
            return Response({"detail": "CV results not found"}, status=404)
        
        cv_data = cv_results[filename]
        roadmap = []
        if "Python" not in cv_data["skills_found"]:
            roadmap.append("Learn Python basics → advanced Python → projects")
        if "SQL" not in cv_data["skills_found"]:
            roadmap.append("Learn SQL → practice databases → data projects")
        if cv_data["experience_years"] < 2:
            roadmap.append("Build 2+ small projects to gain experience")
        roadmap.append("Apply to internships or entry-level positions relevant to your skills")

        return Response({"filename": filename, "skills_found": cv_data["skills_found"], 
                         "experience_years": cv_data["experience_years"], "roadmap": roadmap})


class ExportPDFView(APIView):
    def post(self, request):
        filename = request.data.get("filename")
        if not filename or filename not in uploaded_texts:
            return Response({"detail": "File not found"}, status=404)
        
        content = uploaded_texts[filename]
        pdf_path = os.path.join(EXPORT_DIR, f"{filename}.pdf")
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for line in content.split("\n"):
            pdf.multi_cell(0, 8, line)
        pdf.output(pdf_path)
        return Response({"pdf_path": pdf_path})

class ExportDOCXView(APIView):
    def post(self, request):
        filename = request.data.get("filename")
        if not filename or filename not in uploaded_texts:
            return Response({"detail": "File not found"}, status=404)
        
        content = uploaded_texts[filename]
        docx_path = os.path.join(EXPORT_DIR, f"{filename}.docx")
        doc = Document()
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(docx_path)
        return Response({"docx_path": docx_path})

class ExportPPTXView(APIView):
    def post(self, request):
        filename = request.data.get("filename")
        if not filename or filename not in uploaded_texts:
            return Response({"detail": "File not found"}, status=404)
        
        content = uploaded_texts[filename]
        pptx_path = os.path.join(EXPORT_DIR, f"{filename}.pptx")
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        lines = content.split("\n")
        for i in range(0, len(lines), 5):
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"{filename} - Slide {i//5 + 1}"
            slide.placeholders[1].text = "\n".join(lines[i:i+5])
        prs.save(pptx_path)
        return Response({"pptx_path": pptx_path})


    
