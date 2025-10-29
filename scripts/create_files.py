from docx import Document
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from pymongo import MongoClient
from pptx import Presentation

# --- Connect to MongoDB ---
client = MongoClient("mongodb://localhost:27017/")
db = client["EduNexusAI"]
summaries_collection = db["summaries"]

# --- Fetch data from summaries ---
summaries = list(summaries_collection.find({}, {"_id": 0}))  # exclude _id

# ---------------- WORD EXPORT ----------------
doc = Document()
doc.add_heading("AI Summaries from MongoDB", level=1)

if summaries:
    for s in summaries:
        doc.add_heading(s.get('title', 'Untitled Summary'), level=2)
        doc.add_paragraph(f"Introduction: {s.get('introduction', 'N/A')}")
        doc.add_paragraph(f"Professional Summary: {s.get('professional_summary', 'N/A')}")
        doc.add_paragraph(f"Student Summary: {s.get('student_summary', 'N/A')}")
        doc.add_paragraph(f"Created At: {s.get('created_at', 'N/A')}")
        doc.add_paragraph("")  # spacing
else:
    doc.add_paragraph("No data found in MongoDB.")

doc.save("example.docx")

# ---------------- PDF EXPORT ----------------
pdf_file = "example.pdf"
c = canvas.Canvas(pdf_file, pagesize=A4)
width, height = A4

c.setFont("Helvetica-Bold", 16)
c.drawString(50, height - 50, "AI Summaries from MongoDB")

c.setFont("Helvetica", 12)
y = height - 90

if summaries:
    for s in summaries:
        lines = [
            f"Title: {s.get('title', 'Untitled Summary')}",
            f"Introduction: {s.get('introduction', 'N/A')}",
            f"Professional Summary: {s.get('professional_summary', 'N/A')}",
            f"Student Summary: {s.get('student_summary', 'N/A')}",
            f"Created At: {s.get('created_at', 'N/A')}",
            ""
        ]
        for line in lines:
            c.drawString(50, y, line)
            y -= 20
            if y < 50:
                c.showPage()
                c.setFont("Helvetica", 12)
                y = height - 50
else:
    c.drawString(50, y, "No data found in MongoDB.")

c.save()

# ---------------- POWERPOINT EXPORT ----------------
pptx_file = "example.pptx"
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "AI Summaries from MongoDB"
slide.placeholders[1].text = "Auto-generated using Python"

# Each summary on its own slide
slide_layout = prs.slide_layouts[1]
if summaries:
    for s in summaries:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = s.get('title', 'Untitled Summary')
        content = slide.placeholders[1]
        content.text = (
            f"Introduction:\n{s.get('introduction', 'N/A')}\n\n"
            f"Professional Summary:\n{s.get('professional_summary', 'N/A')}\n\n"
            f"Student Summary:\n{s.get('student_summary', 'N/A')}\n\n"
            f"Created At: {s.get('created_at', 'N/A')}"
        )
else:
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "No Data Found"
    slide.placeholders[1].text = "There are no summaries in MongoDB."

prs.save(pptx_file)

print("example.docx, example.pdf, and example.pptx created successfully!")
